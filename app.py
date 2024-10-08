from flask import Flask, render_template, request
import os
import re
import nltk
from nltk.tokenize import sent_tokenize
from PyPDF2 import PdfReader
import docx
from pptx import Presentation  # Corrected import statement

nltk.download('punkt')  # Download punkt tokenizer if not already downloaded

app = Flask(__name__)

def read_pdf(file):
    reader = PdfReader(file)
    text = ""
    for page in reader.pages:
        text += page.extract_text() + "\n"
    return text

def read_docx(file):
    doc = docx.Document(file)
    return "\n".join([para.text for para in doc.paragraphs])

def read_ppt(file):
    prs = Presentation(file)  # Updated to use the correct import
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def search_answer(text, question):
    sentences = sent_tokenize(text)
    question = question.lower()
    keywords = re.findall(r'\b\w+\b', question)

    best_sentence = None
    best_match_count = 0

    for sentence in sentences:
        sentence_lower = sentence.lower()
        match_count = sum(1 for keyword in keywords if keyword in sentence_lower)

        if match_count > best_match_count:
            best_match_count = match_count
            best_sentence = sentence

    return best_sentence if best_sentence else "I couldn't find an answer. Please try rephrasing your question or upload another file."

def summarize_text(text):
    sentences = sent_tokenize(text)
    if len(sentences) > 5:  # Just a simple way to summarize
        return ' '.join(sentences[:5]) + '...'  # Summarizes to first 5 sentences
    return text  # If text is short, return it as is

@app.route('/', methods=['GET', 'POST'])
def index():
    if request.method == 'POST':
        # Handle file upload and text input
        if 'file' in request.files:
            file = request.files['file']
            file_extension = os.path.splitext(file.filename)[1].lower()

            if file_extension == '.pdf':
                text = read_pdf(file)
            elif file_extension in ['.ppt', '.pptx']:
                text = read_ppt(file)
            elif file_extension in ['.doc', '.docx']:
                text = read_docx(file)
            elif file_extension == '.txt':
                text = file.read().decode('utf-8')
            else:
                return "Unsupported file type. Please upload a PDF, DOCX, PPT, or TXT file."

            question = request.form.get('question')
            action = request.form.get('action')

            if action == 'get_answer' and question:
                answer = search_answer(text, question)
                return render_template('index.html', answer=answer)

            if action == 'summarize':
                summary = summarize_text(text)
                return render_template('index.html', answer=summary)

        text_input = request.form.get('text_input')
        if text_input:
            question = request.form.get('question')
            action = request.form.get('action')

            if action == 'get_answer' and question:
                answer = search_answer(text_input.strip(), question)
                return render_template('index.html', answer=answer)

            if action == 'summarize':
                summary = summarize_text(text_input.strip())
                return render_template('index.html', answer=summary)

    return render_template('index.html')

if __name__ == '__main__':
    app.run(debug=True)
