import os
import torch
from flask import Flask, request, jsonify, render_template
from transformers import BartTokenizer, BartForConditionalGeneration
import nltk
import PyPDF2
from pptx import Presentation
from docx import Document

# Initialize Flask app
app = Flask(__name__)

# Ensure required NLTK resources are downloaded
nltk.download('punkt')

def read_pdf(file_path):
    text = ""
    with open(file_path, 'rb') as file:
        reader = PyPDF2.PdfReader(file)
        for page in reader.pages:
            text += page.extract_text() + "\n"
    return text

def read_ppt(file_path):
    text = ""
    presentation = Presentation(file_path)
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def read_docx(file_path):
    text = ""
    doc = Document(file_path)
    for para in doc.paragraphs:
        text += para.text + "\n"
    return text

def preprocess_text(text):
    # Basic preprocessing
    text = text.replace('\n', ' ').strip()  # Remove newlines and excess whitespace
    return text

def summarize_text(text):
    # Initialize BART model and tokenizer
    tokenizer = BartTokenizer.from_pretrained('facebook/bart-large-cnn')
    model = BartForConditionalGeneration.from_pretrained('facebook/bart-large-cnn')

    # Preprocess text
    text = preprocess_text(text)
    
    # Encode the text and generate summary
    inputs = tokenizer.encode(text, return_tensors='pt', max_length=1024, truncation=True)
    
    # Set a lower maximum length for the summary to ensure conciseness
    max_length = 100  # Adjust this value for a more concise summary
    
    # Generate the summary
    summary_ids = model.generate(
        inputs,
        max_length=max_length,
        num_beams=4,  # Moderate number of beams for focused summaries
        do_sample=False,  # Disable sampling for more deterministic summaries
        early_stopping=True
    )

    # Decode the summary
    summary = tokenizer.decode(summary_ids[0], skip_special_tokens=True)

    return summary

# API to process file and generate summary
@app.route('/summarize', methods=['POST'])
def summarize_file():
    if 'file' not in request.files:
        return jsonify({'error': 'No file part'}), 400

    file = request.files['file']
    if file.filename == '':
        return jsonify({'error': 'No selected file'}), 400

    if file:
        file_path = os.path.join(os.getcwd(), file.filename)
        file.save(file_path)

        file_extension = os.path.splitext(file.filename)[1].lower()
        if file_extension == '.pdf':
            text = read_pdf(file_path)
        elif file_extension in ['.ppt', '.pptx']:
            text = read_ppt(file_path)
        elif file_extension in ['.txt']:
            text = file.read().decode('utf-8')
        elif file_extension in ['.doc', '.docx']:
            text = read_docx(file_path)
        else:
            return jsonify({'error': 'Unsupported file type'}), 400

        # Summarize the text
        summary = summarize_text(text)

        return jsonify({'summary': summary})

@app.route('/')
def index():
    return render_template('index.html')

if __name__ == "__main__":
    app.run(debug=True)
